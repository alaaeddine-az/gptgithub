from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
RED = RGBColor(200, 16, 46)
DARK = RGBColor(40, 40, 40)

def add_title(slide, title, subtitle):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RED

    box2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.6))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(16); p2.font.bold = True; p2.font.color.rgb = DARK

def add_section(slide, top, heading, bullets):
    h = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12), Inches(0.4))
    hp = h.text_frame.paragraphs[0]
    hp.text = heading
    hp.font.size = Pt(16); hp.font.bold = True; hp.font.color.rgb = RED

    b = slide.shapes.add_textbox(Inches(1.0), Inches(top + 0.35), Inches(11.8), Inches(1.35))
    tf = b.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(13); p.font.color.rgb = DARK

def add_footer(slide, text):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(13.33), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    shape.line.color.rgb = RGBColor(220, 220, 220)

    f = slide.shapes.add_textbox(Inches(0.6), Inches(7.02), Inches(12.2), Inches(0.35))
    fp = f.text_frame.paragraphs[0]
    fp.text = f"👉 Message clé : {text}"
    fp.font.size = Pt(13); fp.font.bold = True; fp.font.color.rgb = RED

slides_data = [
    (
        "🔴 SLIDE 1 — ACCULTURATION IA (ENTITÉS MÉTIER)",
        "Acculturer les métiers pour transformer l'IA en levier opérationnel et décisionnel",
        [
            ("Une approche structurée par entité métier", [
                "Ateliers dédiés par domaine (CX, opérations, maintenance, corporate…)",
                "Contenus adaptés aux enjeux spécifiques (cas d'usage RAM)",
                "Alternance vision → cas concrets → mise en pratique",
            ]),
            ("Un format engageant et orienté adoption", [
                "Ateliers interactifs (quizz live, sondages, cas pratiques)",
                "Simulations réelles (prompting, analyse de cas métiers)",
                "Gamification (score individuel / équipe, challenges inter-entités)",
                "Désignation d'AI Champions par entité",
            ]),
        ],
        "Passer d'une compréhension théorique de l'IA à une appropriation concrète par les métiers",
    ),
    (
        "🔴 SLIDE 2 — FORMATION OUTILS (DATA & IA)",
        "Former chaque profil aux bons outils pour industrialiser l'usage de la data et de l'IA",
        [
            ("Parcours de formation différenciés par profil", [
                "Consommateurs (Power BI, dashboards)",
                "Power users (analystes, citizen data)",
                "Experts (data scientists, ML/GenAI engineers)",
            ]),
            ("Trois niveaux de maturité", [
                "Niveau 1 – Usage : Lire, interpréter, interagir avec les outils",
                "Niveau 2 – Autonomie : Construire analyses, dashboards, prompts",
                "Niveau 3 – Industrialisation : Développer modèles, pipelines, agents IA",
            ]),
            ("Modules clés (exemples)", [
                "Data & BI : Power BI, data visualisation, storytelling",
                "Data science : Python, ML, feature engineering",
                "GenAI : prompting, RAG, agents, use cases métier",
                "MLOps / LLMOps : déploiement, monitoring, gouvernance",
            ]),
        ],
        "Créer une chaîne de valeur complète, du consommateur de data au producteur de solutions IA",
    ),
    (
        "🔴 SLIDE 3 — COMMUNITY OF PRACTICE (CoP)",
        "Structurer une communauté interne pour pérenniser l'adoption et accélérer la diffusion de l'IA",
        [
            ("Une organisation en 3 niveaux", [
                "Core team (AI CoE) : Gouvernance, standards, animation",
                "Champions métiers : Relais dans chaque entité",
                "Communauté élargie : Tous les utilisateurs data & IA",
            ]),
            ("Un rôle clé : diffuser, capitaliser, industrialiser", [
                "Partage des bonnes pratiques",
                "Mutualisation des assets (modèles, prompts, use cases)",
                "Accélération de la réutilisation",
                "Support aux équipes projets",
            ]),
            ("Des rituels structurants", [
                "Use case showcases (retours d'expérience)",
                "Tech deep dives (ML, GenAI, MLOps)",
                "Sessions Q&A avec experts",
                "Ateliers inter-entités (cross-fertilisation)",
            ]),
        ],
        "Transformer l'IA en compétence organisationnelle durable, au-delà des projets",
    ),
]

for title, subtitle, sections, footer in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    top = 1.6
    for heading, bullets in sections:
        add_section(slide, top, heading, bullets)
        top += 1.35
    add_footer(slide, footer)

out = "presentation_acculturation_IA_RAM.pptx"
prs.save(out)
print(f"Generated: {out}")
